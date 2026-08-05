#!/usr/bin/env python3
"""Generate the Chinese REIM technical handout and presentation assets.

Outputs
-------
docs/REIM_technical_explanation_zh.docx
docs/REIM_presentation_zh.pptx
docs/assets/*.png

The Markdown source is authoritative. PDF conversion is intentionally handled
by LibreOffice so Chinese fonts remain editable in the DOCX/PPTX sources.
"""

from __future__ import annotations

import re
import textwrap
from datetime import datetime
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image, ImageDraw, ImageFont
import matplotlib

matplotlib.use("Agg")
from matplotlib.font_manager import FontProperties
from matplotlib.mathtext import math_to_image
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor as DocxRGB
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement as PptxOxmlElement
from pptx.util import Inches as PptxInches, Pt as PptxPt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"
MARKDOWN = DOCS / "REIM_technical_explanation_zh.md"
DOCX_OUTPUT = DOCS / "REIM_technical_explanation_zh.docx"
PPTX_OUTPUT = DOCS / "REIM_presentation_zh.pptx"
FONT_FILE = Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc")
FONT = "Noto Sans CJK SC"
MONO_FONT = "Noto Sans Mono CJK SC"


NAVY = "172B3A"
BLUE = "326D9D"
BLUE_LIGHT = "E8F1F8"
ORANGE = "DF7B2F"
ORANGE_LIGHT = "FFF0E6"
TEAL = "3C8F88"
TEAL_LIGHT = "E7F4F2"
GRAY = "667784"
GRAY_LIGHT = "F4F7F9"
GRID = "D7E0E6"
WHITE = "FFFFFF"
BLACK = "18232D"


def hex_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def pil_font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont:
    if bold:
        candidate = Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc")
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.truetype(str(FONT_FILE), size)


def draw_centered_text(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    text: str,
    font: ImageFont.FreeTypeFont,
    fill: str = BLACK,
    spacing: int = 10,
) -> None:
    x0, y0, x1, y1 = box
    wrapped = text
    bbox = draw.multiline_textbbox((0, 0), wrapped, font=font, spacing=spacing, align="center")
    width = bbox[2] - bbox[0]
    height = bbox[3] - bbox[1]
    draw.multiline_text(
        ((x0 + x1 - width) / 2, (y0 + y1 - height) / 2 - bbox[1]),
        wrapped,
        font=font,
        fill=hex_rgb(fill),
        spacing=spacing,
        align="center",
    )


def rounded_box(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    text: str,
    *,
    fill: str,
    outline: str,
    font_size: int = 46,
    radius: int = 24,
    width: int = 5,
) -> None:
    draw.rounded_rectangle(
        box,
        radius=radius,
        fill=hex_rgb(fill),
        outline=hex_rgb(outline),
        width=width,
    )
    draw_centered_text(draw, box, text, pil_font(font_size, bold=True))


def arrow(
    draw: ImageDraw.ImageDraw,
    start: tuple[int, int],
    end: tuple[int, int],
    *,
    color: str = GRAY,
    width: int = 8,
    head: int = 25,
) -> None:
    draw.line([start, end], fill=hex_rgb(color), width=width)
    sx, sy = start
    ex, ey = end
    dx, dy = ex - sx, ey - sy
    length = max((dx * dx + dy * dy) ** 0.5, 1.0)
    ux, uy = dx / length, dy / length
    px, py = -uy, ux
    base_x, base_y = ex - ux * head, ey - uy * head
    points = [
        (ex, ey),
        (base_x + px * head * 0.55, base_y + py * head * 0.55),
        (base_x - px * head * 0.55, base_y - py * head * 0.55),
    ]
    draw.polygon(points, fill=hex_rgb(color))


def save_diagram(image: Image.Image, filename: str) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / filename
    image.save(path, dpi=(220, 220))
    return path


def make_closed_loop_diagram() -> Path:
    image = Image.new("RGB", (2400, 1150), "white")
    d = ImageDraw.Draw(image)
    d.text((70, 45), "REIM 部署闭环：名义控制与选择性恢复", font=pil_font(62, bold=True), fill=hex_rgb(NAVY))

    rounded_box(d, (80, 430, 430, 700), "机器人状态\nsₜ ∈ ℝ²¹", fill=WHITE, outline=TEAL)
    rounded_box(d, (540, 250, 930, 505), "状态版 ACT\n20 步动作块", fill=BLUE_LIGHT, outline=BLUE)
    rounded_box(d, (540, 650, 930, 905), "因果 LSTM\n最近 10 步", fill=ORANGE_LIGHT, outline=ORANGE)

    diamond = [(1090, 450), (1260, 575), (1090, 700), (920, 575)]
    d.polygon(diamond, fill=hex_rgb(ORANGE_LIGHT), outline=hex_rgb(ORANGE))
    d.line(diamond + [diamond[0]], fill=hex_rgb(ORANGE), width=5)
    draw_centered_text(d, (925, 455, 1255, 695), "风险门控\npₜ ≥ 0.20?", pil_font(43, bold=True))

    rounded_box(d, (1380, 220, 1780, 475), "ACT 名义动作", fill=BLUE_LIGHT, outline=BLUE)
    rounded_box(d, (1380, 665, 1780, 920), "恢复策略接管\n成功或 150 步", fill=ORANGE_LIGHT, outline=ORANGE)
    rounded_box(d, (1910, 430, 2300, 700), "执行动作\nMeta-World Sawyer", fill=TEAL_LIGHT, outline=TEAL)

    arrow(d, (430, 505), (540, 380), color=BLUE)
    arrow(d, (430, 625), (540, 780), color=ORANGE)
    arrow(d, (930, 380), (1380, 350), color=BLUE)
    arrow(d, (930, 780), (930, 650), color=ORANGE)
    arrow(d, (930, 780), (930, 650), color=ORANGE)
    arrow(d, (930, 780), (1020, 660), color=ORANGE)
    arrow(d, (1260, 520), (1380, 385), color=BLUE)
    arrow(d, (1260, 630), (1380, 785), color=ORANGE)
    arrow(d, (1780, 350), (1910, 505), color=BLUE)
    arrow(d, (1780, 785), (1910, 625), color=ORANGE)
    arrow(d, (2110, 700), (2110, 1020), color=TEAL)
    arrow(d, (2110, 1020), (250, 1020), color=TEAL)
    arrow(d, (250, 1020), (250, 700), color=TEAL)

    d.text((1280, 455), "否", font=pil_font(34, bold=True), fill=hex_rgb(BLUE))
    d.text((1280, 670), "是", font=pil_font(34, bold=True), fill=hex_rgb(ORANGE))
    d.text((835, 1050), "环境反馈：下一状态 + 任务结果", font=pil_font(34), fill=hex_rgb(TEAL))
    return save_diagram(image, "reim_closed_loop_zh.png")


def make_recovery_curriculum_diagram() -> Path:
    image = Image.new("RGB", (2400, 1060), "white")
    d = ImageDraw.Draw(image)
    d.text((70, 45), "触发对齐恢复课程：从实际接管分布中学习", font=pil_font(62, bold=True), fill=hex_rgb(NAVY))

    boxes = [
        ((80, 340, 430, 690), "扰动 ACT rollout\n+ 因果 LSTM", BLUE_LIGHT, BLUE),
        ((510, 340, 850, 690), "首次风险触发\nτ_collect = 0.10", ORANGE_LIGHT, ORANGE),
        ((930, 280, 1350, 750), "完整 MuJoCo snapshot\nq, q̇, mocap, control\ngoal + task state", TEAL_LIGHT, TEAL),
        ((1430, 180, 1820, 475), "Scripted expert\n从快照继续\n仅保留成功轨迹", WHITE, TEAL),
        ((1430, 565, 1820, 860), "状态—动作对\n42,386 train\n8,212 validation", WHITE, TEAL),
        ((1900, 340, 2320, 690), "确定性恢复 actor\n21→256→256→4\nSmooth L1", ORANGE_LIGHT, ORANGE),
    ]
    for box, text, fill, outline in boxes:
        rounded_box(d, box, text, fill=fill, outline=outline, font_size=39)
    arrow(d, (430, 515), (510, 515), color=ORANGE)
    arrow(d, (850, 515), (930, 515), color=TEAL)
    arrow(d, (1350, 420), (1430, 330), color=TEAL)
    arrow(d, (1350, 610), (1430, 710), color=TEAL)
    arrow(d, (1820, 330), (1900, 460), color=ORANGE)
    arrow(d, (1820, 710), (1900, 585), color=ORANGE)

    d.rounded_rectangle((590, 870, 1810, 990), radius=22, fill=hex_rgb(GRAY_LIGHT), outline=hex_rgb(GRID), width=3)
    draw_centered_text(
        d,
        (600, 875, 1800, 985),
        "训练 bank 与验证 bank 使用不相交 seeds；评测时不调用 expert",
        pil_font(36, bold=True),
        fill=NAVY,
    )
    return save_diagram(image, "recovery_curriculum_zh.png")


def make_state_action_diagram() -> Path:
    image = Image.new("RGB", (2400, 1000), "white")
    d = ImageDraw.Draw(image)
    d.text((70, 45), "状态与动作：当前正式模型使用 21 维语义状态", font=pil_font(62, bold=True), fill=hex_rgb(NAVY))
    labels = [
        ("关节位置", "7D", BLUE),
        ("末端位置+姿态", "7D", BLUE),
        ("物体位置", "3D", TEAL),
        ("目标位置", "3D", TEAL),
        ("夹爪状态", "1D", ORANGE),
    ]
    x, y, h = 80, 280, 235
    widths = [420, 600, 390, 390, 360]
    for (label, dim, color), w in zip(labels, widths):
        d.rounded_rectangle((x, y, x + w, y + h), radius=22, fill=hex_rgb(WHITE), outline=hex_rgb(color), width=5)
        draw_centered_text(d, (x, y + 15, x + w, y + 145), label, pil_font(38, bold=True))
        draw_centered_text(d, (x, y + 120, x + w, y + h - 10), dim, pil_font(42, bold=True), fill=color)
        x += w + 20

    d.text((80, 575), "动作 aₜ ∈ [-1,1]⁴", font=pil_font(48, bold=True), fill=hex_rgb(NAVY))
    action_labels = ["Δx", "Δy", "Δz", "夹爪 g"]
    x = 650
    for idx, label in enumerate(action_labels):
        color = BLUE if idx < 3 else ORANGE
        d.rounded_rectangle((x, 560, x + 350, 790), radius=22, fill=hex_rgb(BLUE_LIGHT if idx < 3 else ORANGE_LIGHT), outline=hex_rgb(color), width=5)
        draw_centered_text(d, (x, 560, x + 350, 790), label, pil_font(58, bold=True), fill=color)
        x += 390
    d.text((80, 880), "边界：state-based 仿真控制，不是视觉端到端或实体机器人 sensing stack", font=pil_font(34), fill=hex_rgb(GRAY))
    return save_diagram(image, "state_action_zh.png")


def make_crn_diagram() -> Path:
    image = Image.new("RGB", (2400, 980), "white")
    d = ImageDraw.Draw(image)
    d.text((70, 45), "公平评测：同一冻结任务库逐 episode 配对", font=pil_font(62, bold=True), fill=hex_rgb(NAVY))
    rounded_box(d, (80, 320, 600, 730), "CRN episode bank\n初态 + 扰动日程\nseed + SHA256\n1,000 episodes", fill=TEAL_LIGHT, outline=TEAL, font_size=42)
    methods = [
        ("ACT", BLUE_LIGHT, BLUE),
        ("ACT + Random Reset", GRAY_LIGHT, GRAY),
        ("ACT + Heuristic Recovery", TEAL_LIGHT, TEAL),
        ("REIM", ORANGE_LIGHT, ORANGE),
    ]
    ys = [170, 360, 550, 740]
    for (label, fill, outline), y in zip(methods, ys):
        rounded_box(d, (1040, y, 1770, y + 150), label, fill=fill, outline=outline, font_size=38)
        arrow(d, (600, 525), (1040, y + 75), color=outline, width=6)
        rounded_box(d, (1920, y, 2320, y + 150), "逐 episode\n成功/步数/干预", fill=WHITE, outline=outline, font_size=31)
        arrow(d, (1770, y + 75), (1920, y + 75), color=outline, width=6)
    return save_diagram(image, "crn_protocol_zh.png")


def make_problem_diagram() -> Path:
    image = Image.new("RGB", (2400, 980), "white")
    d = ImageDraw.Draw(image)
    d.text((70, 45), "一个恢复系统必须同时回答三个问题", font=pil_font(62, bold=True), fill=hex_rgb(NAVY))
    items = [
        ("WHEN", "什么时候离开\n名义 ACT？", "因果 LSTM 风险门控", BLUE_LIGHT, BLUE),
        ("WHERE", "恢复策略应在\n哪些状态学习？", "在线触发的完整 snapshot", TEAL_LIGHT, TEAL),
        ("HOW", "接管后如何完成\n任务并稳定退出？", "持久恢复 option + 清空旧 chunk", ORANGE_LIGHT, ORANGE),
    ]
    for i, (tag, question, answer, fill, outline) in enumerate(items):
        x0 = 80 + i * 770
        x1 = x0 + 690
        d.rounded_rectangle((x0, 235, x1, 835), radius=28, fill=hex_rgb(fill), outline=hex_rgb(outline), width=6)
        draw_centered_text(d, (x0, 250, x1, 390), tag, pil_font(52, bold=True), fill=outline)
        draw_centered_text(d, (x0 + 25, 395, x1 - 25, 610), question, pil_font(44, bold=True))
        d.line((x0 + 70, 640, x1 - 70, 640), fill=hex_rgb(outline), width=4)
        draw_centered_text(d, (x0 + 25, 650, x1 - 25, 815), answer, pil_font(34, bold=True), fill=outline)
    return save_diagram(image, "three_questions_zh.png")


def crop_source_figures() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    framework = Image.open(ROOT / "paper_assets/Figure1_final_framework.png").convert("RGB")
    w, h = framework.size
    framework.crop((0, 0, w, int(h * 0.51))).save(ASSETS / "framework_training_crop.png", quality=95)
    framework.crop((0, int(h * 0.48), w, h)).save(ASSETS / "framework_deployment_crop.png", quality=95)

    result = Image.open(ROOT / "paper_assets/Figure2_final_results.png").convert("RGB")
    w, h = result.size
    cuts = [(0, int(w * 0.335)), (int(w * 0.315), int(w * 0.675)), (int(w * 0.655), w)]
    names = ["result_success_crop.png", "result_intervention_crop.png", "result_robustness_crop.png"]
    for (x0, x1), name in zip(cuts, names):
        result.crop((x0, 0, x1, int(h * 0.93))).save(ASSETS / name, quality=95)


def build_visual_assets() -> None:
    crop_source_figures()
    make_closed_loop_diagram()
    make_recovery_curriculum_diagram()
    make_state_action_diagram()
    make_crn_diagram()
    make_problem_diagram()


# ---------------------------------------------------------------------------
# DOCX renderer
# ---------------------------------------------------------------------------


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_border(cell, color: str = GRID, size: str = "6") -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        tag = f"w:{edge}"
        element = borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), size)
        element.set(qn("w:color"), color)


def set_run_font_docx(run, size: float | None = None, *, bold: bool | None = None, color: str | None = None, mono: bool = False) -> None:
    name = MONO_FONT if mono else FONT
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold
    if color is not None:
        run.font.color.rgb = DocxRGB(*hex_rgb(color))


INLINE_PATTERN = re.compile(r"(\*\*.+?\*\*|`.+?`)")
INLINE_MATH_PATTERN = re.compile(r"\\\((.+?)\\\)")


def readable_inline_math(text: str) -> str:
    """Convert the small inline-math subset used by the handout to Unicode."""

    superscript = str.maketrans("-0123456789", "⁻⁰¹²³⁴⁵⁶⁷⁸⁹")

    def convert(match: re.Match[str]) -> str:
        value = match.group(1)
        value = value.replace(r"\beta", "β").replace(r"\tau", "τ").replace(r"\times", "×")
        value = re.sub(r"10\^\{(-?\d+)\}", lambda item: "10" + item.group(1).translate(superscript), value)
        value = re.sub(r"_\{([^{}]+)\}", r"_\1", value)
        value = value.replace(r"\,", " ").replace(r"\;", " ")
        return value

    return INLINE_MATH_PATTERN.sub(convert, text)


def add_inline_runs(paragraph, text: str, *, size: float = 10.5, color: str = BLACK) -> None:
    text = readable_inline_math(text)
    pos = 0
    for match in INLINE_PATTERN.finditer(text):
        if match.start() > pos:
            run = paragraph.add_run(text[pos : match.start()])
            set_run_font_docx(run, size, color=color)
        token = match.group(0)
        if token.startswith("**"):
            run = paragraph.add_run(token[2:-2])
            set_run_font_docx(run, size, bold=True, color=color)
        else:
            run = paragraph.add_run(token[1:-1])
            set_run_font_docx(run, size, color=NAVY, mono=True)
        pos = match.end()
    if pos < len(text):
        run = paragraph.add_run(text[pos:])
        set_run_font_docx(run, size, color=color)


def add_page_number(paragraph) -> None:
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = " PAGE "
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "end")
    run._r.extend([fld_char1, instr_text, fld_char2])
    set_run_font_docx(run, 9, color=GRAY)


def configure_doc_styles(doc: Document) -> None:
    normal = doc.styles["Normal"]
    normal.font.name = FONT
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
    normal.font.size = Pt(10.5)
    normal.paragraph_format.space_after = Pt(5)
    normal.paragraph_format.line_spacing = 1.18

    for style_name, size, color in (
        ("Title", 26, NAVY),
        ("Heading 1", 18, NAVY),
        ("Heading 2", 14, BLUE),
        ("Heading 3", 11.5, TEAL),
    ):
        style = doc.styles[style_name]
        style.font.name = FONT
        style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = DocxRGB(*hex_rgb(color))
        style.paragraph_format.keep_with_next = True
        style.paragraph_format.space_before = Pt(12 if style_name != "Title" else 0)
        style.paragraph_format.space_after = Pt(6)

    for style_name in ("List Bullet", "List Number"):
        style = doc.styles[style_name]
        style.font.name = FONT
        style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
        style.font.size = Pt(10.5)

    if "Code Block" not in [s.name for s in doc.styles]:
        style = doc.styles.add_style("Code Block", WD_STYLE_TYPE.PARAGRAPH)
        style.font.name = MONO_FONT
        style._element.rPr.rFonts.set(qn("w:eastAsia"), MONO_FONT)
        style.font.size = Pt(8.5)
        style.paragraph_format.left_indent = Cm(0.6)
        style.paragraph_format.right_indent = Cm(0.6)
        style.paragraph_format.space_before = Pt(4)
        style.paragraph_format.space_after = Pt(6)


def add_doc_title_page(doc: Document) -> None:
    for _ in range(2):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("REIM 项目技术精讲")
    set_run_font_docx(run, 30, bold=True, color=NAVY)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Recovery-Enhanced Imitation Learning\nfor Robust Embodied Robot Manipulation")
    set_run_font_docx(run, 16, bold=True, color=BLUE)
    doc.add_paragraph()
    image_path = ROOT / "results/figures/recovery_operation_sequence_frames/08_reim_success_seed8300042_t062.png"
    if image_path.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(image_path), width=Cm(14.8))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_inline_runs(p, "Meta-World PickPlace · MuJoCo Sawyer · ACT + Causal LSTM + Recovery Imitation", size=11, color=GRAY)
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_inline_runs(p, "技术说明、实验口径、代码地图与答辩 FAQ\n2026-08-04", size=11, color=NAVY)
    p.add_run().add_break(WD_BREAK.PAGE)


def add_manual_toc(doc: Document, markdown_text: str) -> None:
    p = doc.add_paragraph("目录", style="Title")
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    headings = []
    for line in markdown_text.splitlines():
        if line.startswith("## "):
            headings.append(line[3:].strip())
    for heading in headings:
        p = doc.add_paragraph(style="List Number")
        add_inline_runs(p, heading, size=10.5, color=NAVY)
    p = doc.add_paragraph()
    p.add_run().add_break(WD_BREAK.PAGE)


def add_code_block(doc: Document, lines: Sequence[str]) -> None:
    p = doc.add_paragraph(style="Code Block")
    p.paragraph_format.keep_together = True
    p_pr = p._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), GRAY_LIGHT)
    p_pr.append(shd)
    run = p.add_run("\n".join(lines))
    set_run_font_docx(run, 8.5, color=BLACK, mono=True)


def add_markdown_table(doc: Document, lines: Sequence[str]) -> None:
    parsed = []
    for line in lines:
        cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
        parsed.append(cells)
    if len(parsed) >= 2 and all(re.fullmatch(r":?-{3,}:?", c.replace(" ", "")) for c in parsed[1]):
        parsed.pop(1)
    if not parsed:
        return
    columns = max(len(row) for row in parsed)
    table = doc.add_table(rows=len(parsed), cols=columns)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    table.autofit = True
    for r_idx, row in enumerate(parsed):
        for c_idx in range(columns):
            cell = table.cell(r_idx, c_idx)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            set_cell_border(cell)
            if r_idx == 0:
                set_cell_shading(cell, NAVY)
            elif r_idx % 2 == 0:
                set_cell_shading(cell, GRAY_LIGHT)
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT if c_idx == 0 else WD_ALIGN_PARAGRAPH.CENTER
            text = row[c_idx] if c_idx < len(row) else ""
            add_inline_runs(p, text, size=8.5, color=WHITE if r_idx == 0 else BLACK)
            for run in p.runs:
                run.bold = r_idx == 0
    doc.add_paragraph().paragraph_format.space_after = Pt(2)


IMAGE_PATTERN = re.compile(r"!\[(.*?)\]\((.*?)\)")


def resolve_markdown_image(raw_path: str) -> Path:
    path = (MARKDOWN.parent / raw_path).resolve()
    return path


def add_markdown_image(doc: Document, caption: str, raw_path: str) -> None:
    path = resolve_markdown_image(raw_path)
    if not path.exists():
        p = doc.add_paragraph()
        add_inline_runs(p, f"[缺失图片：{path}]", size=9, color=ORANGE)
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    with Image.open(path) as im:
        ratio = im.width / max(im.height, 1)
    width = Cm(16.4)
    if ratio < 1.0:
        width = Cm(10.2)
    elif ratio > 2.4:
        width = Cm(17.2)
    p.add_run().add_picture(str(path), width=width)
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(caption)
    set_run_font_docx(run, 8.5, color=GRAY)
    cap.paragraph_format.keep_with_next = False


def add_quote(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.7)
    p.paragraph_format.right_indent = Cm(0.5)
    p_pr = p._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), ORANGE_LIGHT)
    p_pr.append(shd)
    add_inline_runs(p, text, size=10, color=NAVY)


def normalize_display_math(source: str) -> str:
    """Normalize valid LaTeX that is slightly stricter than mathtext syntax."""

    source = " ".join(part.strip() for part in source.splitlines() if part.strip())
    source = re.sub(
        r"\\text\{([^}]*)\}",
        lambda match: r"\mathrm{" + match.group(1).replace(" ", r"\;") + "}",
        source,
    )
    source = re.sub(r"\\mathcal\s+([A-Za-z])", r"\\mathcal{\1}", source)
    source = re.sub(r"\\mathbf\s+([A-Za-z0-9])", r"\\mathbf{\1}", source)
    return source


def add_display_equation(doc: Document, source: str, number: int) -> None:
    """Render a display equation as a high-resolution, centered math image."""

    ASSETS.mkdir(parents=True, exist_ok=True)
    output = ASSETS / f"equation_{number:02d}.png"
    equation = normalize_display_math(source)
    equation_dpi = 220
    math_to_image(
        f"${equation}$",
        output,
        prop=FontProperties(size=17),
        dpi=equation_dpi,
        format="png",
        color=f"#{BLACK}",
    )
    with Image.open(output) as image:
        natural_width_cm = image.width / equation_dpi * 2.54
    display_width_cm = min(15.6, max(2.2, natural_width_cm))

    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.paragraph_format.keep_together = True
    paragraph.paragraph_format.space_before = Pt(2)
    paragraph.paragraph_format.space_after = Pt(4)
    paragraph.add_run().add_picture(str(output), width=Cm(display_width_cm))


def render_markdown_to_docx() -> Path:
    text = MARKDOWN.read_text(encoding="utf-8")
    doc = Document()
    section = doc.sections[0]
    section.page_width = Cm(21.0)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(1.8)
    section.bottom_margin = Cm(1.7)
    section.left_margin = Cm(1.9)
    section.right_margin = Cm(1.9)
    section.header_distance = Cm(0.7)
    section.footer_distance = Cm(0.7)
    configure_doc_styles(doc)

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = header.add_run("REIM · 项目技术精讲")
    set_run_font_docx(run, 8.5, color=GRAY)
    add_page_number(section.footer.paragraphs[0])

    add_doc_title_page(doc)
    add_manual_toc(doc, text)

    lines = text.splitlines()
    index = 0
    equation_number = 0
    in_code = False
    code_lines: list[str] = []
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code:
                add_code_block(doc, code_lines)
                code_lines = []
                in_code = False
            else:
                in_code = True
            index += 1
            continue
        if in_code:
            code_lines.append(line)
            index += 1
            continue
        if stripped == r"\[":
            equation_lines: list[str] = []
            index += 1
            while index < len(lines) and lines[index].strip() != r"\]":
                equation_lines.append(lines[index])
                index += 1
            if index >= len(lines):
                raise ValueError("Unterminated display equation in Markdown source")
            equation_number += 1
            add_display_equation(doc, "\n".join(equation_lines), equation_number)
            index += 1
            continue
        if stripped.startswith("|") and "|" in stripped[1:]:
            table_lines = []
            while index < len(lines) and lines[index].strip().startswith("|"):
                table_lines.append(lines[index])
                index += 1
            add_markdown_table(doc, table_lines)
            continue
        image_match = IMAGE_PATTERN.fullmatch(stripped)
        if image_match:
            add_markdown_image(doc, image_match.group(1), image_match.group(2))
            index += 1
            continue
        if stripped.startswith("# "):
            # The source title is already represented by the title page.
            index += 1
            continue
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=1)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=2)
        elif stripped.startswith("#### "):
            doc.add_heading(stripped[5:], level=3)
        elif stripped.startswith("> "):
            add_quote(doc, stripped[2:])
        elif re.match(r"^\d+\.\s+", stripped):
            content = re.sub(r"^\d+\.\s+", "", stripped)
            p = doc.add_paragraph(style="List Number")
            add_inline_runs(p, content)
        elif stripped.startswith("- "):
            p = doc.add_paragraph(style="List Bullet")
            add_inline_runs(p, stripped[2:])
        elif stripped == "---":
            p = doc.add_paragraph()
            p_pr = p._p.get_or_add_pPr()
            borders = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "8")
            bottom.set(qn("w:color"), GRID)
            borders.append(bottom)
            p_pr.append(borders)
        elif stripped:
            p = doc.add_paragraph()
            if stripped.startswith("\\[") or stripped.startswith("\\]") or stripped.startswith("+"):
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            add_inline_runs(p, stripped)
        index += 1

    doc.core_properties.title = "REIM 项目技术精讲"
    doc.core_properties.subject = "Meta-World PickPlace failure recovery with ACT, causal LSTM and trigger-aligned recovery imitation"
    doc.core_properties.author = "REIM Project"
    doc.core_properties.last_modified_by = "REIM Project"
    doc.core_properties.keywords = "REIM, ACT, failure recovery, Meta-World, Sawyer, imitation learning"
    doc.core_properties.created = datetime(2026, 8, 4)
    doc.core_properties.modified = datetime(2026, 8, 4)
    doc.save(DOCX_OUTPUT)
    return DOCX_OUTPUT


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------


SLIDE_W = PptxInches(13.333)
SLIDE_H = PptxInches(7.5)


def rgb(value: str) -> RGBColor:
    return RGBColor(*hex_rgb(value))


def set_pptx_run_font(run, *, size: float, bold: bool = False, color: str = BLACK, name: str = FONT) -> None:
    run.font.name = name
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "zh-CN")


def set_shape_fill(shape, fill: str, line: str | None = None, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line or fill)


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 18,
    bold: bool = False,
    color: str = BLACK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
    fill: str | None = None,
    line: str | None = None,
    radius: bool = False,
):
    if radius or fill is not None:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        set_shape_fill(shape, fill or WHITE, line or fill or WHITE)
        frame = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        frame = shape.text_frame
    frame.clear()
    frame.margin_left = PptxInches(margin)
    frame.margin_right = PptxInches(margin)
    frame.margin_top = PptxInches(margin)
    frame.margin_bottom = PptxInches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = PptxPt(0)
    run = p.add_run()
    run.text = text
    set_pptx_run_font(run, size=size, bold=bold, color=color)
    return shape


def add_bullets(slide, x, y, w, h, bullets: Sequence[str], *, size: float = 17, color: str = BLACK, accent: str = ORANGE) -> None:
    box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = PptxInches(0.05)
    frame.margin_right = PptxInches(0.04)
    frame.margin_top = 0
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.level = 0
        p.text = item
        p.font.size = PptxPt(size)
        p.font.name = FONT
        p.font.color.rgb = rgb(color)
        p.space_after = PptxPt(8)
        p.line_spacing = 1.08
        p._p.get_or_add_pPr().insert(0, PptxOxmlElement("a:buChar"))
        p._p.get_or_add_pPr()[0].set("char", "•")
    return box


def add_slide_base(prs: Presentation, title: str, number: int, core: str, *, section: str = "REIM"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(WHITE)
    add_textbox(slide, 0.55, 0.25, 10.8, 0.52, title, size=27, bold=True, color=NAVY)
    add_textbox(slide, 11.5, 0.31, 1.15, 0.28, section, size=10, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(0.55), PptxInches(0.89), PptxInches(12.2), PptxInches(0.025))
    set_shape_fill(line, GRID)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptxInches(0.55), PptxInches(6.83), PptxInches(11.65), PptxInches(0.44))
    set_shape_fill(strip, GRAY_LIGHT, GRID)
    frame = strip.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "本页结论  " + core
    set_pptx_run_font(run, size=11.2, bold=True, color=NAVY)
    add_textbox(slide, 12.35, 6.91, 0.4, 0.2, str(number), size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
    return slide


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, *, border: str | None = None):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    if ratio >= box_ratio:
        draw_w = w
        draw_h = w / ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    pic = slide.shapes.add_picture(str(path), PptxInches(draw_x), PptxInches(draw_y), width=PptxInches(draw_w), height=PptxInches(draw_h))
    if border:
        pic.line.color.rgb = rgb(border)
        pic.line.width = PptxPt(1.5)
    return pic


def add_metric_card(slide, x, y, w, h, value, label, *, color=ORANGE, fill=ORANGE_LIGHT, value_size=31):
    shape = add_textbox(slide, x, y, w, h, "", fill=fill, line=color, radius=True)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value; set_pptx_run_font(r, size=value_size, bold=True, color=color)
    p2 = frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = label; set_pptx_run_font(r, size=11.5, bold=True, color=NAVY)
    return shape


def add_notes(slide, text: str) -> None:
    try:
        frame = slide.notes_slide.notes_text_frame
        frame.text = text
    except Exception:
        pass


def add_simple_table(slide, x, y, w, h, rows: Sequence[Sequence[str]], widths: Sequence[float] | None = None, *, font_size=13.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
    table = table_shape.table
    if widths:
        for col, width in zip(table.columns, widths):
            col.width = PptxInches(width)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.margin_left = PptxInches(0.05)
            cell.margin_right = PptxInches(0.05)
            cell.margin_top = PptxInches(0.03)
            cell.margin_bottom = PptxInches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = rgb(NAVY)
            elif r_idx == len(rows) - 1:
                cell.fill.fore_color.rgb = rgb(ORANGE_LIGHT)
            elif r_idx % 2 == 0:
                cell.fill.fore_color.rgb = rgb(GRAY_LIGHT)
            else:
                cell.fill.fore_color.rgb = rgb(WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            for run in p.runs:
                set_pptx_run_font(run, size=font_size, bold=(r_idx == 0 or r_idx == len(rows) - 1), color=WHITE if r_idx == 0 else (ORANGE if r_idx == len(rows) - 1 else BLACK))
    return table_shape


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    bg_frame = ROOT / "results/figures/recovery_operation_sequence_frames/08_reim_success_seed8300042_t062.png"
    add_picture_contain(slide, bg_frame, 7.0, 0, 6.333, 7.5)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, PptxInches(7.7), SLIDE_H)
    set_shape_fill(overlay, NAVY, NAVY)
    add_textbox(slide, 0.72, 1.0, 6.2, 1.25, "REIM", size=54, bold=True, color=WHITE)
    add_textbox(slide, 0.75, 2.1, 5.9, 1.3, "面向机器人操作失败的\n恢复增强模仿学习", size=28, bold=True, color=WHITE)
    add_textbox(slide, 0.78, 3.65, 5.7, 0.8, "ACT · 因果风险门控 · 触发对齐恢复", size=17, bold=True, color="A9DAD6")
    add_textbox(slide, 0.78, 5.45, 5.6, 0.55, "Meta-World PickPlace / MuJoCo Sawyer", size=14, color="DDE8EE")
    add_textbox(slide, 0.78, 6.35, 5.6, 0.35, "技术精讲与实验复现 · 2026-08-04", size=11, color="AFC0CA")
    add_notes(slide, "开场只讲一句：REIM 让模仿策略在执行偏离后能够识别风险并切换到专门的恢复控制。强调右图是 MuJoCo 仿真帧。")

    # 2. Task
    slide = add_slide_base(prs, "任务：Sawyer PickPlace 的抓取、搬运与放置", 2, "在固定任务语义下，研究执行扰动引起的闭环失败。", section="任务")
    add_picture_contain(slide, ROOT / "results/figures/recovery_operation_sequence_frames/01_act_initial_seed8300042_t000.png", 0.65, 1.2, 4.3, 4.85, border=GRID)
    add_picture_contain(slide, ASSETS / "state_action_zh.png", 5.1, 1.25, 7.6, 3.25)
    add_bullets(slide, 5.35, 4.52, 7.0, 1.95, [
        "历史任务 PickPlace-v2；Meta-World 3.1.1 实际注册为 pick-place-v3",
        "21 维语义状态 → 4 维笛卡尔增量与夹爪动作",
        "最大 200 步；成功由 Meta-World goal-reaching 信号判定",
        "状态输入仿真研究，不是视觉策略或实体 Sawyer 实验",
    ], size=15.2)
    add_notes(slide, "先把边界讲清：任务没变，只是 Meta-World API 从历史 v2 迁移到维护中的 v3。当前输入含精确物体和目标位置。")

    # 3. Failure
    slide = add_slide_base(prs, "为什么强 ACT 仍然会失败？", 3, "示范拟合准确不等于扰动下闭环鲁棒。", section="问题")
    frames = [
        ("物体位移 t=3", "02_act_disturbance_seed8300042_t003.png"),
        ("最佳进展仍未抬升 t=68", "03_act_unrecovered_seed8300042_t068.png"),
        ("ACT 超时 t=200", "04_act_failure_seed8300042_t200.png"),
    ]
    for i, (label, name) in enumerate(frames):
        x = 0.65 + i * 4.15
        add_picture_contain(slide, ROOT / "results/figures/recovery_operation_sequence_frames" / name, x, 1.35, 3.8, 3.8, border=ORANGE if i else GRID)
        add_textbox(slide, x, 5.15, 3.8, 0.35, label, size=13.5, bold=True, color=ORANGE if i else NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, 5.75, 11.5, 0.65, "小误差 → 新状态偏离专家分布 → 下一步动作更差 → 误差沿时间累积", size=21, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=ORANGE_LIGHT, line=ORANGE, radius=True)
    add_notes(slide, "不要把问题说成 ACT 太弱。ACT 已经是强名义策略，失败来自模仿学习闭环中的 covariate shift 与 compounding error。")

    # 4. Three questions
    slide = add_slide_base(prs, "研究问题：恢复系统必须同时回答三件事", 4, "关键不只是学恢复动作，还包括切换时机与训练分布。", section="问题")
    add_picture_contain(slide, ASSETS / "three_questions_zh.png", 0.65, 1.15, 12.0, 5.45)
    add_notes(slide, "三个问题分别对应 LSTM 门控、触发快照数据和持久恢复 option。为后续模块建立一一对应关系。")

    # 5. Overview
    slide = add_slide_base(prs, "REIM 一句话方案：选择性恢复闭环", 5, "ACT 负责常规操作，LSTM 判断风险，恢复策略只在需要时接管。", section="方法")
    add_picture_contain(slide, ASSETS / "reim_closed_loop_zh.png", 0.6, 1.15, 12.1, 5.5)
    add_notes(slide, "部署闭环是全文唯一主线。先讲低风险走 ACT，高风险走恢复，再讲环境反馈。")

    # 6. ACT
    slide = add_slide_base(prs, "模块一：状态版 ACT 名义策略", 6, "动作块表达短期操作结构，时间集成减少逐步控制抖动。", section="ACT")
    add_textbox(slide, 0.7, 1.2, 7.5, 0.8, "sₜ → CVAE Transformer → 20-action chunk → temporal ensemble → aₜ", size=19, bold=True, color=BLUE, align=PP_ALIGN.CENTER, fill=BLUE_LIGHT, line=BLUE, radius=True)
    add_bullets(slide, 0.8, 2.2, 7.25, 3.85, [
        "500 条成功专家轨迹，25,969 个 transition",
        "训练时 posterior encoder 看当前状态与未来专家动作块",
        "L1 重建 + β·KL，β=10；推理固定 z=0",
        "chunk=20，hidden=256，8 heads，encoder/decoder=3/4 层",
        "重叠 chunk 按 exp(-0.05·i) 加权；切换时清空历史",
    ], size=16.2, color=BLACK)
    add_metric_card(slide, 8.6, 1.4, 1.75, 1.35, "20", "动作块长度", color=BLUE, fill=BLUE_LIGHT)
    add_metric_card(slide, 10.55, 1.4, 1.75, 1.35, "6.63M", "参数量", color=BLUE, fill=BLUE_LIGHT)
    add_metric_card(slide, 8.6, 3.05, 1.75, 1.35, "0.004013", "验证动作 L1", color=TEAL, fill=TEAL_LIGHT, value_size=23)
    add_metric_card(slide, 10.55, 3.05, 1.75, 1.35, "100%", "示范成功率", color=TEAL, fill=TEAL_LIGHT)
    add_textbox(slide, 8.65, 4.8, 3.65, 1.15, "state-based ACT\n不声称复现原始多相机版本", size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "ACT 的价值是强名义技能，而不是失败恢复。CVAE encoder 训练时使用，推理时固定 z=0，保证确定性。")

    # 7. Failure data
    slide = add_slide_base(prs, "构造风险数据：让 ACT 在扰动下暴露真实失效状态", 7, "监测器训练于学习策略实际诱导的轨迹，而非静态异常样本。", section="风险数据")
    add_picture_contain(slide, ROOT / "results/figures/recovery_operation_sequence_frames/02_act_disturbance_seed8300042_t003.png", 0.75, 1.25, 4.1, 4.65, border=ORANGE)
    add_bullets(slide, 5.1, 1.3, 7.2, 2.3, [
        "执行 2,000 个扰动 ACT episode",
        "动作高斯噪声 σ=0.15",
        "观测噪声 σ=0.01",
        "物体位移概率 0.02、幅度 0.04 m",
    ], size=17)
    add_metric_card(slide, 5.25, 4.0, 2.1, 1.45, "285,443", "因果时序窗口", color=ORANGE, fill=ORANGE_LIGHT, value_size=25)
    add_metric_card(slide, 7.65, 4.0, 2.1, 1.45, "39.65%", "episode 成功率", color=ORANGE, fill=ORANGE_LIGHT, value_size=25)
    add_metric_card(slide, 10.05, 4.0, 2.1, 1.45, "64.09%", "正标签占比", color=ORANGE, fill=ORANGE_LIGHT, value_size=25)
    add_textbox(slide, 5.3, 5.75, 6.7, 0.55, "失败事件：掉落 / 越界 / 抓取失败 / 停滞 / 偏离 / 超时", size=15.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "数据来自冻结 ACT 的在线 rollout。进度规则只在发生物体交互后开启，避免把正常接近阶段误标成停滞。")

    # 8. Detector
    slide = add_slide_base(prs, "模块二：最近 10 步的因果 LSTM 风险门控", 8, "输入只到当前时刻；标签判断事件是否位于 [t,t+10]。", section="风险门控")
    add_textbox(slide, 0.8, 1.3, 6.25, 0.85, "sₜ₋₉:ₜ (10×21) → LSTM 128 → MLP 64 → sigmoid → pₜ", size=19, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, fill=ORANGE_LIGHT, line=ORANGE, radius=True)
    add_bullets(slide, 0.9, 2.45, 6.0, 3.6, [
        "严格因果窗口：不使用 t 之后的状态",
        "早期窗口右侧补零，并传入有效长度",
        "Weighted BCEWithLogitsLoss",
        "trajectory-grouped 80%/20% 划分",
        "冻结部署阈值 τ_on=0.20",
    ], size=16.5)
    add_metric_card(slide, 7.65, 1.45, 1.85, 1.3, "91.8%", "Precision", color=ORANGE, fill=ORANGE_LIGHT)
    add_metric_card(slide, 9.75, 1.45, 1.85, 1.3, "92.3%", "Recall", color=ORANGE, fill=ORANGE_LIGHT)
    add_metric_card(slide, 8.7, 3.05, 1.85, 1.3, "92.1%", "窗口 F1", color=ORANGE, fill=ORANGE_LIGHT)
    add_textbox(slide, 7.55, 4.75, 4.15, 1.05, "称为“因果风险监测器”\n而不是纯长时域预测器", size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "重点解释 inclusive label：输入严格因果，但标签区间包含当前事件，所以不能把 92.1% F1 等同于提前 10 步预测。")

    # 9. Recovery challenge
    slide = add_slide_base(prs, "恢复学习的真正难点：状态分布必须对齐", 9, "相同物体位置不代表相同动力学与控制器状态。", section="恢复学习")
    add_textbox(slide, 0.8, 1.35, 5.2, 4.9, "近似人工 reset\n\n只移动 object xyz\n\n缺少 q̇ / gripper / mocap\ncontrol / task bookkeeping\n\n训练—部署分布错位", size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRAY, radius=True)
    add_textbox(slide, 7.3, 1.35, 5.2, 4.9, "在线触发 exact snapshot\n\n保存完整 MuJoCo 动态状态\n\n来自 ACT + LSTM\n实际会到达的接管时刻\n\n训练分布 ≈ 部署触发分布", size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER, fill=TEAL_LIGHT, line=TEAL, radius=True)
    add_textbox(slide, 6.1, 3.1, 1.1, 0.7, "→", size=42, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_notes(slide, "这是方法创新的核心动机。人工重置只控制几何位置，而恢复行为还依赖速度、夹爪和任务内部状态。")

    # 10. Curriculum
    slide = add_slide_base(prs, "核心创新：触发对齐的恢复课程", 10, "在第一次风险触发时保存完整状态，再从同一状态收集成功纠正轨迹。", section="恢复学习")
    add_picture_contain(slide, ASSETS / "recovery_curriculum_zh.png", 0.6, 1.1, 12.15, 5.55)
    add_notes(slide, "采集阈值 0.10 比部署阈值 0.20 更宽松，是为了覆盖更多仍可恢复的边缘状态。评测阶段绝不调用 expert。")

    # 11. Recovery actor
    slide = add_slide_base(prs, "模块三：轻量确定性恢复 actor", 11, "恢复能力来自触发状态上的监督纠正模仿，而非额外在线试错。", section="恢复策略")
    add_textbox(slide, 0.8, 1.35, 6.0, 0.95, "state 21 → tanh 256 → tanh 256 → action 4 → clip[-1,1]", size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER, fill=TEAL_LIGHT, line=TEAL, radius=True)
    add_bullets(slide, 0.9, 2.55, 6.0, 3.45, [
        "42,386 个训练状态—动作对",
        "8,212 个轨迹不相交验证对",
        "Smooth L1，40 epochs，batch 512，lr=10⁻³",
        "输入高斯增强 σ=0.005",
        "验证 Smooth L1 = 0.006391",
    ], size=16.4)
    add_metric_card(slide, 7.6, 1.45, 2.1, 1.4, "72,452", "actor 参数", color=TEAL, fill=TEAL_LIGHT, value_size=27)
    add_metric_card(slide, 10.0, 1.45, 2.1, 1.4, "0", "RL 环境训练步", color=ORANGE, fill=ORANGE_LIGHT, value_size=31)
    add_metric_card(slide, 7.6, 3.25, 2.1, 1.4, "0", "policy-gradient 更新", color=ORANGE, fill=ORANGE_LIGHT, value_size=31)
    add_metric_card(slide, 10.0, 3.25, 2.1, 1.4, "60,598", "导出一致性状态", color=TEAL, fill=TEAL_LIGHT, value_size=25)
    add_textbox(slide, 7.65, 5.15, 4.35, 0.78, "最终正式算法：supervised recovery imitation\n不是 PPO / RL", size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "这里必须明确口径。正式 checkpoint 无 critic、optimizer 或 rollout state，动作导出审计误差为零。")

    # 12. Arbitration
    slide = add_slide_base(prs, "部署仲裁：一旦触发，恢复 option 持续接管", 12, "持久恢复避免在未回到安全状态时频繁切换。", section="闭环控制")
    add_picture_contain(slide, ASSETS / "reim_closed_loop_zh.png", 0.7, 1.2, 7.1, 4.55)
    rows = [
        ["参数", "冻结值"],
        ["触发阈值", "0.20"],
        ["恢复预算", "150 步"],
        ["最小接管", "150 步"],
        ["清除窗口", "200 步"],
        ["中途风险释放", "禁用"],
    ]
    add_simple_table(slide, 8.15, 1.45, 4.15, 3.75, rows, widths=[2.2, 1.95], font_size=14)
    add_textbox(slide, 8.2, 5.45, 4.05, 0.65, "所有 hand-off 清空 ACT temporal ensemble", size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, fill=ORANGE_LIGHT, line=ORANGE, radius=True)
    add_notes(slide, "min_steps 等于 budget、clear_steps 大于 budget，因此恢复中途不会根据风险清除退出。预算用满且任务未结束时才返回 ACT，并清空动作块。")

    # 13. CRN
    slide = add_slide_base(prs, "实验协议：冻结 CRN 任务库逐 episode 配对", 13, "同一初态、同一扰动日程，让差异主要来自控制器。", section="实验")
    add_picture_contain(slide, ASSETS / "crn_protocol_zh.png", 0.65, 1.15, 8.05, 5.3)
    add_bullets(slide, 8.85, 1.55, 3.65, 3.7, [
        "主实验：20% 扰动，1,000 episodes/方法",
        "鲁棒性：5 个等级，各 200 episodes",
        "Success：Wilson 95% CI",
        "差异与步数：episode bootstrap",
        "episode bank 与原始记录均有 SHA256",
    ], size=14.4)
    add_textbox(slide, 8.95, 5.45, 3.45, 0.65, "配对设计比独立随机抽样\n更直接地衡量“救回/损害”", size=14.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "Random Reset 的第二次尝试有独立保留 seed，并在表里明确标注。其他方法消费相同 episode specification。")

    # 14. Baselines
    slide = add_slide_base(prs, "四种方法逐层回答：恢复是否有效，门控是否必要？", 14, "比较覆盖无恢复、重试、语义门控恢复和完整 REIM。", section="实验")
    cards = [
        ("ACT", "无干预\n名义能力", BLUE_LIGHT, BLUE),
        ("ACT + Random Reset", "失败后一次\nfresh-task retry", GRAY_LIGHT, GRAY),
        ("ACT + Heuristic Recovery", "语义规则触发\n同一恢复 actor", TEAL_LIGHT, TEAL),
        ("REIM", "因果 LSTM 触发\n同一恢复 actor", ORANGE_LIGHT, ORANGE),
    ]
    for i, (name, desc, fill, outline) in enumerate(cards):
        x = 0.65 + i * 3.15
        add_textbox(slide, x, 1.45, 2.75, 3.75, "", fill=fill, line=outline, radius=True)
        add_textbox(slide, x + 0.15, 1.8, 2.45, 0.9, name, size=18 if i != 2 else 15.5, bold=True, color=outline, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, 3.0, 2.35, 1.2, desc, size=15.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.25, 4.45, 2.25, 0.4, ["基准", "额外机会", "恢复作用", "门控价值"][i], size=11.5, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.4, 5.65, 10.5, 0.55, "Heuristic 与 REIM 使用完全相同的恢复 actor → 差异主要来自仲裁规则", size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "Random Reset 的 outcome 不是 recovery rate；Heuristic 使用模拟器级语义信息，是强但 privileged 的门控对照。")

    # 15. Main success
    slide = add_slide_base(prs, "主要结果：REIM 将成功率从 73.4% 提升到 90.4%", 15, "恢复闭环在 20% 扰动下带来 +17.0 个百分点。", section="结果")
    add_picture_contain(slide, ASSETS / "result_success_crop.png", 0.65, 1.15, 8.25, 5.25)
    add_metric_card(slide, 9.25, 1.55, 2.5, 1.65, "90.4%", "REIM task success", color=ORANGE, fill=ORANGE_LIGHT, value_size=37)
    add_metric_card(slide, 9.25, 3.55, 2.5, 1.65, "+17.0 pp", "相对 ACT", color=TEAL, fill=TEAL_LIGHT, value_size=31)
    add_textbox(slide, 9.1, 5.65, 2.8, 0.55, "平均步数 69.3\nACT 为 93.3", size=14.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_notes(slide, "先讲绝对成功率，再讲配对提升。不要在这一页塞入所有 p 值。")

    # 16. Pairing/intervention
    slide = add_slide_base(prs, "REIM 不靠更多干预：成功更高，干预更少", 16, "相比启发式门控，REIM 少干预 76 个 episode。", section="结果")
    add_picture_contain(slide, ASSETS / "result_intervention_crop.png", 0.65, 1.25, 7.5, 4.9)
    add_metric_card(slide, 8.45, 1.35, 1.75, 1.35, "170", "ACT 失败被救回", color=TEAL, fill=TEAL_LIGHT)
    add_metric_card(slide, 10.45, 1.35, 1.75, 1.35, "0", "ACT 成功被损害", color=TEAL, fill=TEAL_LIGHT)
    add_metric_card(slide, 8.45, 3.0, 1.75, 1.35, "+2.7 pp", "vs. Heuristic", color=ORANGE, fill=ORANGE_LIGHT, value_size=25)
    add_metric_card(slide, 10.45, 3.0, 1.75, 1.35, "−10.4%", "相对干预负担", color=ORANGE, fill=ORANGE_LIGHT, value_size=25)
    add_textbox(slide, 8.4, 4.8, 3.85, 1.0, "REIM：65.5% episodes 干预\nHeuristic：73.1%", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "这页回答常见质疑：REIM 不是因为更频繁地启动恢复才成功，实际干预率更低。")

    # 17. Robustness
    slide = add_slide_base(prs, "扰动越强，恢复闭环的价值越明显", 17, "40% 扰动下优势扩大到 +32.0 个百分点。", section="结果")
    add_picture_contain(slide, ASSETS / "result_robustness_crop.png", 0.65, 1.1, 8.8, 5.35)
    rows = [
        ["噪声", "ACT", "REIM"],
        ["0%", "100.0", "100.0"],
        ["10%", "97.0", "99.5"],
        ["20%", "76.5", "92.0"],
        ["30%", "53.5", "79.5"],
        ["40%", "31.5", "63.5"],
    ]
    add_simple_table(slide, 9.55, 1.4, 2.9, 4.3, rows, widths=[0.8, 1.0, 1.1], font_size=12.5)
    add_textbox(slide, 9.65, 5.9, 2.7, 0.45, "单位：task success (%)", size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    add_notes(slide, "0% 时两者都为 100%，说明恢复机制没有牺牲名义任务；噪声增加时收益快速扩大。")

    # 18. Ablation
    slide = add_slide_base(prs, "组件消融：大提升来自恢复，进一步提升来自学习式门控", 18, "相同恢复 actor 下，LSTM 同时提高成功率并减少干预。", section="消融")
    add_picture_contain(slide, ROOT / "paper_assets/Figure3_final_ablation.png", 0.75, 1.2, 7.2, 5.15)
    add_bullets(slide, 8.25, 1.5, 4.0, 3.2, [
        "ACT，无恢复：73.4%",
        "Heuristic gate + recovery：87.7%",
        "LSTM gate + recovery：90.4%",
        "干预率：73.1% → 65.5%",
    ], size=16.5)
    add_textbox(slide, 8.25, 4.85, 4.0, 1.05, "证据支持两层结论\n恢复动作有效 + 学习门控更优", size=17.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, fill=GRAY_LIGHT, line=GRID, radius=True)
    add_notes(slide, "不把耦合的 snapshot collection 和 recovery imitation 强行拆成两项独立贡献，实验没有单独隔离它们。")

    # 19. Gate sensitivity
    slide = add_slide_base(prs, "门控阈值：在相同干预负担下，LSTM 仍显著更好", 19, "τ=0.175 与启发式干预率几乎相同，但成功率高 6.5 pp。", section="门控诊断")
    add_picture_contain(slide, ROOT / "paper_assets/Figure4_gate_sensitivity.png", 0.65, 1.1, 8.2, 5.45)
    add_metric_card(slide, 9.15, 1.35, 2.6, 1.4, "92.5% vs 86.0%", "matched-burden success", color=ORANGE, fill=ORANGE_LIGHT, value_size=23)
    add_metric_card(slide, 9.15, 3.05, 2.6, 1.4, "76.5% vs 76.0%", "intervention rate", color=TEAL, fill=TEAL_LIGHT, value_size=23)
    add_metric_card(slide, 9.15, 4.75, 2.6, 1.25, "p = 0.0044", "exact paired test", color=NAVY, fill=GRAY_LIGHT, value_size=24)
    add_notes(slide, "说明这是冻结之后的独立 sensitivity bank，不用于回头选择主阈值。主部署点仍是 0.20。")

    # 20. Detector diagnostics
    slide = add_slide_base(prs, "风险监测器：闭环有效，但提前预测能力有边界", 20, "标准窗口 F1 很高；严格未来事件 F1 为 60.0%。", section="诚实诊断")
    add_picture_contain(slide, ROOT / "paper_assets/Figure3_detector.png", 0.75, 1.2, 5.15, 5.1)
    rows = [
        ["诊断", "结果"],
        ["Inclusive Precision", "91.8%"],
        ["Inclusive Recall", "92.3%"],
        ["Inclusive F1", "92.1%"],
        ["当前事件正窗口", "86.1%"],
        ["Strict future F1", "60.0%"],
        ["事件前轨迹告警", "191/258 = 74.0%"],
        ["中位提前量", "1 step"],
    ]
    add_simple_table(slide, 6.4, 1.35, 5.75, 4.7, rows, widths=[3.2, 2.55], font_size=13.2)
    add_textbox(slide, 6.45, 6.15, 5.65, 0.35, "结论依赖端到端闭环结果，而不是单独包装分类指标", size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_notes(slide, "主动解释 86.1% offset=0 能提高可信度。称 causal risk monitor，不称 long-horizon failure forecaster。")

    # 21. Operation
    slide = add_slide_base(prs, "一条完整恢复轨迹：ACT 超时，REIM 在 62 步完成", 21, "风险在轨迹尚可修复时触发，恢复 actor 完成后续抓取与搬运。", section="定性案例")
    add_picture_contain(slide, ROOT / "paper_assets/Figure5_operation_sequence.png", 0.6, 1.02, 12.15, 5.7)
    add_notes(slide, "这是独立定性 seed 8300042：t=3 物体位移 0.048m，t=9 风险 0.259，53 个恢复动作，t=62 成功。明确不是实体照片。")

    # 22. Conclusion
    slide = add_slide_base(prs, "结论：REIM 实现选择性早期纠正，而非万能失败修复", 22, "触发分布对齐 + 因果门控 + 持久恢复构成一个可审计闭环。", section="总结")
    add_picture_contain(slide, ASSETS / "reim_closed_loop_zh.png", 0.65, 1.2, 7.25, 4.2)
    add_metric_card(slide, 8.25, 1.35, 1.8, 1.35, "90.4%", "任务成功率", color=ORANGE, fill=ORANGE_LIGHT)
    add_metric_card(slide, 10.25, 1.35, 1.8, 1.35, "+17.0 pp", "相对 ACT", color=TEAL, fill=TEAL_LIGHT, value_size=25)
    add_metric_card(slide, 9.25, 3.0, 1.8, 1.35, "−10.4%", "vs. heuristic burden", color=BLUE, fill=BLUE_LIGHT, value_size=25)
    add_bullets(slide, 8.25, 4.75, 4.0, 1.45, [
        "边界：单任务、21D privileged state、无实体实验",
        "约 99% 触发发生在首次 lift 前",
        "下一步：post-drop、多任务、视觉与真实 Sawyer",
    ], size=13.5)
    add_notes(slide, "最后主动限定证据边界。项目当前最强结论是 selective early correction and task recovery under simulated execution disturbances。")

    prs.core_properties.title = "REIM 项目技术精讲"
    prs.core_properties.subject = "ACT, causal LSTM and trigger-aligned recovery imitation on Meta-World PickPlace"
    prs.core_properties.author = "REIM Project"
    prs.core_properties.last_modified_by = "REIM Project"
    prs.core_properties.keywords = "REIM, ACT, Meta-World, Sawyer, failure recovery"
    prs.core_properties.created = datetime(2026, 8, 4)
    prs.core_properties.modified = datetime(2026, 8, 4)
    prs.save(PPTX_OUTPUT)
    return PPTX_OUTPUT


def main() -> None:
    if not MARKDOWN.is_file():
        raise FileNotFoundError(MARKDOWN)
    build_visual_assets()
    docx_path = render_markdown_to_docx()
    pptx_path = build_presentation()
    print(f"docx={docx_path}")
    print(f"pptx={pptx_path}")
    print(f"assets={ASSETS}")


if __name__ == "__main__":
    main()
